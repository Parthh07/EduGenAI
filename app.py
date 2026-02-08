import os
from flask import Flask, render_template, request, jsonify
import google.generativeai as genai
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from dotenv import load_dotenv

load_dotenv()

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

api_key = os.getenv("GEMINI_API_KEY")
genai.configure(api_key=api_key)
model = genai.GenerativeModel('gemini-flash-latest')

current_context = ""

# --- FILE EXTRACTORS ---
def get_pdf_text(filepath):
    reader = PdfReader(filepath)
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""
    return text

def get_docx_text(filepath):
    doc = Document(filepath)
    text = []
    for paragraph in doc.paragraphs:
        text.append(paragraph.text)
    return "\n".join(text)

def get_pptx_text(filepath):
    prs = Presentation(filepath)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def get_txt_text(filepath):
    with open(filepath, 'r', encoding='utf-8') as f:
        return f.read()

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    global current_context
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
        
    filepath = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
    file.save(filepath)
    
    try:
        ext = os.path.splitext(filepath)[1].lower()
        if ext == '.pdf': text = get_pdf_text(filepath)
        elif ext == '.docx': text = get_docx_text(filepath)
        elif ext == '.pptx': text = get_pptx_text(filepath)
        elif ext == '.txt': text = get_txt_text(filepath)
        else: return jsonify({'error': 'Unsupported file format'}), 400

        current_context = text[:30000]
        return jsonify({'message': 'Document Processed Successfully!'})
        
    except Exception as e:
        return jsonify({'error': f"Error reading file: {str(e)}"}), 500

@app.route('/upload_text', methods=['POST'])
def upload_text():
    global current_context
    data = request.get_json()
    text = data.get('text', '').strip()
    if not text: return jsonify({'error': 'Canvas is empty!'}), 400
    current_context = text[:30000]
    return jsonify({'message': 'Text Processed Successfully!'})

@app.route('/generate_notes', methods=['POST'])
def generate_notes():
    global current_context
    if not current_context:
        return jsonify({'error': '⚠️ Canvas Empty. Please Upload document first.'}), 400
    
    data = request.get_json() or {}
    length_setting = data.get('length', 'Standard')
    
    length_instructions = {
        "Brief": "summary (bullet points only, max 300 words)",
        "Standard": "standard notes with definitions, key concepts, and summary",
        "Detailed": "deep-dive tutorial with examples and long explanations"
    }
    instruction = length_instructions.get(length_setting, length_instructions['Standard'])

    try:
        prompt = f"""
        Act as an expert tutor. Analyze the text below.
        GOAL: Create {instruction}.
        Format output in clean HTML (<h3>, <ul>, <li>, <strong>, <p>). Do NOT use Markdown.
        Text: {current_context[:15000]}
        """
        response = model.generate_content(prompt)
        clean_text = response.text.replace('```html', '').replace('```', '')
        return jsonify({'content': clean_text})
    except Exception as e:
        return jsonify({'error': f"AI Error: {str(e)}"}), 500

# --- UPDATED QUIZ GENERATOR ---
@app.route('/generate_quiz', methods=['POST'])
def generate_quiz():
    if not current_context:
        return jsonify({'error': '⚠️ Upload document first.'}), 400

    data = request.get_json() or {}
    q_type = data.get('type', 'mcq') # 'mcq' or 'theory'
    difficulty = data.get('difficulty', 'Medium')
    
    # Logic for Theory Questions (2, 6, 10 Marks)
    if q_type == 'theory':
        marks = data.get('marks', '2') # 2, 6, or 10
        
        prompt = f"""
        Generate 3 theoretical questions based on the text for a {marks}-marks exam category.
        
        **Criteria:**
        - If 2 Marks: Short definition or concept check (1-2 sentences).
        - If 6 Marks: Explain a concept with points/steps (paragraph + bullets).
        - If 10 Marks: Long detailed answer, case study, or "Explain in detail" type.

        Return strictly raw JSON format:
        [ {{"question": "Question text...", "answer": "Model Answer..."}} ]
        
        Do not use Markdown.
        Text: {current_context[:15000]}
        """
        
    # Logic for MCQs
    else:
        prompt = f"""
        Generate 5 Multiple Choice Questions (MCQs).
        **Difficulty:** {difficulty}
        Return strictly raw JSON format:
        [ {{"question": "...", "options": ["A", "B", "C", "D"], "answer": "Option"}} ]
        Do not use Markdown.
        Text: {current_context[:15000]}
        """

    try:
        response = model.generate_content(prompt)
        clean_text = response.text.replace('```json', '').replace('```', '')
        return jsonify({'content': clean_text, 'type': q_type})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    app.run(debug=True)